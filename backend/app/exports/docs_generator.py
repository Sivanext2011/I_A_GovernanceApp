"""Generate document exports with live data from processed datasets."""

from pathlib import Path
from io import BytesIO
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.drawing.image import Image as XLImage
from openpyxl.utils import get_column_letter
import plotly.graph_objects as go

from app.config import settings
from app.analytics import (
    compute_kpis, compute_monthly_trend, compute_department_comparison,
    compute_downloads_vs_reuse, compute_leaderboard,
)
from app.services.data_service import data_store
from app.utils import get_logger
from app.utils.departments import TEAMS

logger = get_logger(__name__)

HEADER_FONT = Font(bold=True, color="FFFFFF", size=9)
HEADER_FILL = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
TEAM_FILL = PatternFill(start_color="2E86AB", end_color="2E86AB", fill_type="solid")
THIN_BORDER = Border(
    left=Side(style="thin"), right=Side(style="thin"),
    top=Side(style="thin"), bottom=Side(style="thin"),
)

MONTH_NAMES = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
               "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
TEAMS_ORDER = ["Billing", "Charging", "SDC CS&DFE", "SDC Billing&MW"]
YEAR = "2026"


def _get_2026_months() -> list[str]:
    """Get only 2026 months from available data."""
    return sorted([m for m in data_store.get_available_months() if m.startswith(YEAR)])


def _month_label(canonical: str) -> str:
    try:
        return MONTH_NAMES[int(canonical.split("-")[1]) - 1]
    except Exception:
        return canonical


def _generate_chart_png(chart_type: str, months_2026: list[str]) -> Path:
    """Generate chart as PNG file and return path."""
    team = "Overall"
    filepath = settings.EXPORT_DIR / f"_doc_chart_{chart_type}.png"

    # Filter trend data to only 2026 months
    fig = go.Figure()

    if chart_type == "savings_trend":
        labels = [_month_label(m) for m in months_2026]
        total, auto, reuse = [], [], []
        for m in months_2026:
            k = compute_kpis([m], team)
            total.append(k["total_savings"])
            auto.append(k["automation_savings"])
            reuse.append(k["asset_savings"])
        fig.add_trace(go.Bar(x=labels, y=total, name="Total Savings", marker_color="#1F4E79"))
        fig.add_trace(go.Bar(x=labels, y=auto, name="Automation", marker_color="#14b8a6"))
        fig.add_trace(go.Bar(x=labels, y=reuse, name="Reuse", marker_color="#f59e0b"))
        fig.update_layout(title="Monthly Savings Trend (Overall) - 2026", barmode="group")

    elif chart_type == "savings_pct":
        labels = [_month_label(m) for m in months_2026]
        pcts = []
        for m in months_2026:
            k = compute_kpis([m], team)
            pcts.append(k["savings_percent"] or 0)
        fig.add_trace(go.Bar(x=labels, y=pcts, name="Savings %", marker_color="#8b5cf6"))
        fig.update_layout(title="Savings % Trend (Overall) - 2026", yaxis_title="%")

    elif chart_type == "dept_comparison":
        comp = compute_department_comparison(months_2026)
        fig.add_trace(go.Bar(x=comp["teams"], y=comp["total_savings"], name="Total Savings", marker_color="#3b82f6"))
        fig.add_trace(go.Scatter(
            x=comp["teams"], y=comp["savings_percent"], name="Savings %",
            yaxis="y2", mode="lines+markers", line=dict(color="#ef4444", width=3)
        ))
        fig.update_layout(
            title="Department wise Savings - 2026",
            yaxis2=dict(overlaying="y", side="right", title="%"),
        )

    elif chart_type == "downloads_vs_reuse":
        labels = [_month_label(m) for m in months_2026]
        downloads, reuse = [], []
        for m in months_2026:
            k = compute_kpis([m], team)
            downloads.append(k["total_downloads"])
            reuse.append(k["total_reused_with_savings"])
        fig.add_trace(go.Bar(x=labels, y=downloads, name="Downloads", marker_color="#3b82f6"))
        fig.add_trace(go.Bar(x=labels, y=reuse, name="Reused with Savings", marker_color="#10b981"))
        fig.update_layout(title="Downloads vs Reuse (Overall) - 2026", barmode="group")

    fig.update_layout(template="plotly_white", font=dict(size=11), margin=dict(t=50, b=40, l=50, r=50))
    fig.write_image(str(filepath), width=900, height=400, scale=2)
    return filepath


# ============================================================
# EXCEL GENERATION
# ============================================================

def generate_monthly_savings_report() -> Path:
    """Generate Monthly Savings Report Excel with data sheet and charts sheet."""
    output_path = settings.EXPORT_DIR / "Monthly_Savings_Report.xlsx"
    months_2026 = _get_2026_months()
    all_months_canonical = [f"{YEAR}-{str(i).zfill(2)}" for i in range(1, 13)]

    wb = Workbook()

    # === SHEET 1: Monthly Data ===
    ws = wb.active
    ws.title = f"Monthly Data {YEAR}"
    _write_monthly_data_sheet(ws, months_2026, all_months_canonical)

    # === SHEET 2: Charts ===
    ws_charts = wb.create_sheet("Charts")
    _write_charts_sheet(ws_charts, months_2026)

    wb.save(output_path)
    logger.info(f"Monthly Savings Report generated: {output_path}")
    return output_path


def _write_monthly_data_sheet(ws, months_2026, all_months_canonical):
    """Write the monthly KPI data table."""
    # ROW 1: Team group headers
    ws.cell(row=1, column=1, value="")
    ws.cell(row=1, column=2, value="")
    col = 3
    for team in TEAMS_ORDER:
        ws.merge_cells(start_row=1, start_column=col, end_row=1, end_column=col + 7)
        cell = ws.cell(row=1, column=col, value=team)
        cell.font = Font(bold=True, color="FFFFFF", size=10)
        cell.fill = TEAM_FILL
        cell.alignment = Alignment(horizontal="center")
        col += 8
    ws.merge_cells(start_row=1, start_column=col, end_row=1, end_column=col + 2)
    cell = ws.cell(row=1, column=col, value="Overall")
    cell.font = Font(bold=True, color="FFFFFF", size=10)
    cell.fill = TEAM_FILL
    cell.alignment = Alignment(horizontal="center")

    # ROW 2: Column headers
    headers = ["Month", "Monetization KPI"]
    for team in TEAMS_ORDER:
        headers += [
            f"Total Number of Assets Download - {team}",
            f"Total Number of Assets Reused with Savings - {team}",
            f"Pending Feedback - {team}",
            f"Asset Savings - {team}",
            f"Automation Savings - {team}",
            f"Total Savings - {team}",
            f"Billability Hours - {team}",
            f"{team} Savings %",
        ]
    headers += ["Total Savings - Monetization", "Billability Hours - Monetization", "Monetization Savings %"]

    for c_idx, h in enumerate(headers, 1):
        cell = ws.cell(row=2, column=c_idx, value=h)
        cell.font = HEADER_FONT
        cell.fill = HEADER_FILL
        cell.alignment = Alignment(horizontal="center", wrap_text=True)
        cell.border = THIN_BORDER

    # DATA ROWS: Jan-Dec
    for r_idx, month_canonical in enumerate(all_months_canonical, 3):
        month_label = _month_label(month_canonical)
        has_data = month_canonical in months_2026

        ws.cell(row=r_idx, column=1, value=month_label).border = THIN_BORDER
        ws.cell(row=r_idx, column=2, value="16.00%").border = THIN_BORDER

        col = 3
        overall_total_savings = 0
        overall_billability = 0

        for team in TEAMS_ORDER:
            if has_data:
                kpis = compute_kpis([month_canonical], team)
                vals = [
                    kpis["total_downloads"] or "",
                    kpis["total_reused_with_savings"] or "",
                    kpis["pending_feedback"] or "",
                    kpis["asset_savings"] or "",
                    kpis["automation_savings"] or "",
                    kpis["total_savings"],
                    kpis["billability_hours"] or "",
                    f"{kpis['savings_percent']:.0f}%" if kpis["savings_percent"] else "0%",
                ]
                overall_total_savings += kpis["total_savings"]
                overall_billability += kpis["billability_hours"]
            else:
                vals = ["", "", "", "", "", 0, "", ""]

            for val in vals:
                cell = ws.cell(row=r_idx, column=col, value=val)
                cell.border = THIN_BORDER
                cell.alignment = Alignment(horizontal="right")
                col += 1

        if has_data:
            overall_pct = (overall_total_savings / overall_billability * 100) if overall_billability > 0 else None
            ws.cell(row=r_idx, column=col, value=round(overall_total_savings, 2)).border = THIN_BORDER
            ws.cell(row=r_idx, column=col + 1, value=round(overall_billability, 2)).border = THIN_BORDER
            ws.cell(row=r_idx, column=col + 2, value=f"{overall_pct:.2f}%" if overall_pct else "#DIV/0!").border = THIN_BORDER
        else:
            ws.cell(row=r_idx, column=col, value=0).border = THIN_BORDER
            ws.cell(row=r_idx, column=col + 1, value=0).border = THIN_BORDER
            ws.cell(row=r_idx, column=col + 2, value="#DIV/0!").border = THIN_BORDER

    # YTD ROW
    ytd_row = 15
    ws.cell(row=ytd_row, column=1, value="YTD").border = THIN_BORDER
    ws.cell(row=ytd_row, column=1).font = Font(bold=True)
    ws.cell(row=ytd_row, column=2, value="16.00%").border = THIN_BORDER

    col = 3
    ytd_overall_savings = 0
    ytd_overall_billability = 0

    for team in TEAMS_ORDER:
        kpis = compute_kpis(months_2026, team) if months_2026 else {
            "total_downloads": 0, "total_reused_with_savings": 0, "pending_feedback": 0,
            "asset_savings": 0, "automation_savings": 0, "total_savings": 0,
            "billability_hours": 0, "savings_percent": None,
        }
        vals = [
            kpis["total_downloads"] or "",
            kpis["total_reused_with_savings"] or "",
            kpis["pending_feedback"] or "",
            kpis["asset_savings"] or "",
            kpis["automation_savings"] or "",
            kpis["total_savings"],
            kpis["billability_hours"] or "",
            f"{kpis['savings_percent']:.0f}%" if kpis["savings_percent"] else "0%",
        ]
        ytd_overall_savings += kpis["total_savings"]
        ytd_overall_billability += kpis["billability_hours"]

        for val in vals:
            cell = ws.cell(row=ytd_row, column=col, value=val)
            cell.border = THIN_BORDER
            cell.alignment = Alignment(horizontal="right")
            cell.font = Font(bold=True)
            col += 1

    ytd_pct = (ytd_overall_savings / ytd_overall_billability * 100) if ytd_overall_billability > 0 else None
    ws.cell(row=ytd_row, column=col, value=round(ytd_overall_savings, 2)).border = THIN_BORDER
    ws.cell(row=ytd_row, column=col + 1, value=round(ytd_overall_billability, 2)).border = THIN_BORDER
    ws.cell(row=ytd_row, column=col + 2, value=f"{ytd_pct:.2f}%" if ytd_pct else "0.00%").border = THIN_BORDER
    for c in range(col, col + 3):
        ws.cell(row=ytd_row, column=c).font = Font(bold=True)

    # Auto-fit columns
    for col_idx in range(1, ws.max_column + 1):
        max_len = 0
        for row_idx in range(1, ws.max_row + 1):
            cell = ws.cell(row=row_idx, column=col_idx)
            max_len = max(max_len, len(str(cell.value or "")))
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 2, 35)


def _write_charts_sheet(ws, months_2026):
    """Write Monetization Analytics charts as images into the Charts sheet."""
    ws.cell(row=1, column=1, value="Monetization Analytics - Overall (2026)")
    ws.cell(row=1, column=1).font = Font(bold=True, size=14)

    chart_types = ["savings_trend", "savings_pct", "dept_comparison", "downloads_vs_reuse"]
    chart_titles = ["Monthly Savings Trend", "Savings % Trend", "Department wise Savings", "Downloads vs Reuse"]

    row = 3
    for chart_type, title in zip(chart_types, chart_titles):
        try:
            chart_path = _generate_chart_png(chart_type, months_2026)
            ws.cell(row=row, column=1, value=title).font = Font(bold=True, size=11)
            row += 1
            img = XLImage(str(chart_path))
            img.width = 700
            img.height = 310
            ws.add_image(img, f"A{row}")
            row += 18  # space for image
        except Exception as e:
            ws.cell(row=row, column=1, value=f"Chart generation failed: {e}")
            row += 2


# ============================================================
# POWERPOINT GENERATION
# ============================================================

def generate_asset_presentation(period: str = "monthly") -> Path:
    """Generate Asset Presentation PowerPoint with charts, top reusers, and YTD."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    output_path = settings.EXPORT_DIR / f"Asset_{period.replace('-', '_').title()}.pptx"
    months_2026 = _get_2026_months()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_text_box(slide, "Monetization Analytics - 2026", Inches(1), Inches(2.5),
                  Inches(11), Inches(1.5), size=36, bold=True, color="1F4E79")
    _add_text_box(slide, f"Period: {period.title()} | Overall", Inches(1), Inches(4),
                  Inches(11), Inches(0.8), size=18, color="475569")

    # --- Slide 2: YTD KPIs ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, "Year-to-Date KPIs (Overall)", Inches(0.5), Inches(0.3),
                  Inches(12), Inches(0.7), size=22, bold=True, color="1F4E79")
    ytd_kpis = compute_kpis(months_2026, "Overall")
    kpi_data = [
        ["Total Savings", "Savings %", "Downloads", "Reused", "Pending Feedback", "Billability Hours"],
        [
            f"{ytd_kpis['total_savings']:,.2f}",
            f"{ytd_kpis['savings_percent']:.2f}%" if ytd_kpis["savings_percent"] else "N/A",
            str(ytd_kpis["total_downloads"]),
            str(ytd_kpis["total_reused_with_savings"]),
            str(ytd_kpis["pending_feedback"] or 0),
            f"{ytd_kpis['billability_hours']:,.2f}",
        ]
    ]
    _add_table(slide, kpi_data, Inches(0.5), Inches(1.2), Inches(12), Inches(1.2))

    # Team-wise YTD
    _add_text_box(slide, "Team-wise YTD", Inches(0.5), Inches(3), Inches(12), Inches(0.6),
                  size=16, bold=True, color="2E86AB")
    team_data = [["Team", "Total Savings", "Savings %", "Downloads", "Reused", "Pending", "Billability"]]
    for t in TEAMS_ORDER:
        k = compute_kpis(months_2026, t)
        team_data.append([
            t, f"{k['total_savings']:,.2f}",
            f"{k['savings_percent']:.2f}%" if k["savings_percent"] else "N/A",
            str(k["total_downloads"]), str(k["total_reused_with_savings"]),
            str(k["pending_feedback"] or 0), f"{k['billability_hours']:,.2f}",
        ])
    _add_table(slide, team_data, Inches(0.5), Inches(3.8), Inches(12), Inches(2.5))

    # --- Slide 3: Monthly Savings Trend ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, "Monthly Savings Trend", Inches(0.5), Inches(0.2),
                  Inches(12), Inches(0.6), size=20, bold=True, color="1F4E79")
    chart_path = _generate_chart_png("savings_trend", months_2026)
    slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1), Inches(12), Inches(5.5))

    # --- Slide 4: Savings % Trend ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, "Savings % Trend", Inches(0.5), Inches(0.2),
                  Inches(12), Inches(0.6), size=20, bold=True, color="1F4E79")
    chart_path = _generate_chart_png("savings_pct", months_2026)
    slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1), Inches(12), Inches(5.5))

    # --- Slide 5: Department wise Savings ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, "Department wise Savings", Inches(0.5), Inches(0.2),
                  Inches(12), Inches(0.6), size=20, bold=True, color="1F4E79")
    chart_path = _generate_chart_png("dept_comparison", months_2026)
    slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1), Inches(12), Inches(5.5))

    # --- Slide 6: Downloads vs Reuse ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, "Downloads vs Reuse", Inches(0.5), Inches(0.2),
                  Inches(12), Inches(0.6), size=20, bold=True, color="1F4E79")
    chart_path = _generate_chart_png("downloads_vs_reuse", months_2026)
    slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1), Inches(12), Inches(5.5))

    # --- Slide 7: Top Asset Reusers per Department ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, "Top Asset Reusers by Department (YTD 2026)", Inches(0.5), Inches(0.2),
                  Inches(12), Inches(0.6), size=20, bold=True, color="1F4E79")

    lb_data = [["#", "Name", "Department", "Reuse Saving", "Automation Saving", "Total Savings"]]
    for t in TEAMS_ORDER:
        lb = compute_leaderboard(months_2026, t, top_n=3)
        for i, entry in enumerate(lb, 1):
            lb_data.append([
                str(i), entry["name"], t,
                f"{entry['reuse_saving']:,.2f}",
                f"{entry['automation_saving']:,.2f}",
                f"{entry['total_savings']:,.2f}",
            ])
    _add_table(slide, lb_data, Inches(0.5), Inches(1), Inches(12), Inches(5.5))

    prs.save(str(output_path))

    # Cleanup temp chart images
    for f in settings.EXPORT_DIR.glob("_doc_chart_*.png"):
        f.unlink(missing_ok=True)

    logger.info(f"Asset Presentation generated: {output_path}")
    return output_path


def _add_text_box(slide, text, left, top, width, height, size=14, bold=False, color="000000"):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor.from_string(color)


def _add_table(slide, data, left, top, width, height):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    rows = len(data)
    cols = len(data[0]) if data else 1
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
